from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Colours — clean, neutral
BG     = RGBColor(0xFF, 0xFE, 0xFC)   # near-white
INK    = RGBColor(0x1A, 0x22, 0x2E)   # near-black
MUTED  = RGBColor(0x55, 0x65, 0x78)   # grey
ACCENT = RGBColor(0x1A, 0x6B, 0x5E)   # dark teal

W, H = Inches(13.33), Inches(7.5)     # 16:9 widescreen

slides = [
    {
        "title": "Parallax",
        "bullets": [
            "Retirement planning simulator built for financial advisors",
            "98 years of real market data. Path-consistent Monte Carlo.",
            "Three account types: Taxable, Traditional, Roth",
            "Built by an advisor. Not a software company.",
        ],
        "note": "Cover"
    },
    {
        "title": "The problem with every other tool",
        "bullets": [
            "Built around fear — outputs are calibrated to sell conservative products",
            "Isolated stats: a success rate with no context for what's driving it",
            "Levers don't interact — changing one input doesn't ripple through the plan",
            "Advisors can't safely show the upside. The tool has a bias; it just hides it.",
            "Clients are managed by anxiety rather than an accurate picture of their situation",
        ],
    },
    {
        "title": "What Parallax does differently",
        "bullets": [
            "No thesis. Equally willing to say 'you're fine, spend more' as 'this won't last'",
            "The same gauge reads green when news is good and red when it isn't",
            "Models interactions — every lever moves every other lever in real time",
            "The advisor points at a number. The client sees exactly why it moved.",
            "Neutrality is the product. That's why you can trust it in front of a client.",
        ],
    },
    {
        "title": "The engine",
        "bullets": [
            "98 years of real annual returns (1928–2025), not synthetic data",
            "Block-bootstrap Monte Carlo — resamples real historical blocks, not random draws",
            "Path-consistent: every year of a simulation uses the same return sequence",
            "Models: sequence-of-returns risk, longevity, healthcare cost trajectory, LTC",
            "Tax: withdrawal sequencing across Taxable / Traditional / Roth, tax multiplier lever",
            "Accumulation phase + pension income + social security + one-time spending events",
            "Pure computation — no UI, no defaults that quietly bias the output",
        ],
    },
    {
        "title": "The tax picture",
        "bullets": [
            "Three account types modeled independently and correctly",
            "Withdrawal sequencing: which bucket gets drawn first, and what that costs in taxes",
            "Roth conversion scenarios: run the same household with and without conversion",
            "Tax multiplier: stress-test the plan against higher future ordinary income rates",
            "After-tax vs pre-tax balances are never conflated in the output",
            "The gap: no tool today shows the RETIREMENT cost of a tax decision in real time",
        ],
    },
    {
        "title": "What advisors see — Scenarios tab",
        "bullets": [
            "Baseline column + 1–2 alternative scenarios side by side",
            "Each column: Monte Carlo success % + delta vs baseline the moment you run",
            "Levers: Retirement Age, SS Start Age, Spending, One-Time Event, Allocation, Savings",
            "Hero chart: expected wealth path in today's dollars, current age to plan end",
            "Deterministic line — identical inputs produce identical lines, always",
            "Allocation change fans the lines apart by exactly the compounding difference",
            "Aggressive allocation → higher expected line, lower success circle. That trade-off is the point.",
        ],
    },
    {
        "title": "What advisors see — Sequencing tab",
        "bullets": [
            "Same household plan, run through real historical market sequences",
            "Lines: retire into 1929, 1966 (lost decade), 1973 (stagflation), 2000, 2008",
            "Each line = the same spending, same allocation, different market order",
            "Shows sequence-of-returns risk directly — no abstraction, no jargon",
            "Outcome card per line: first-decade return, lowest balance, survival age or ran-dry age",
            "Advisor can say: 'Here's your plan in the worst sequence of the last 100 years'",
        ],
    },
    {
        "title": "Who it's for",
        "bullets": [
            "Fee-only RIAs, personal-CFO model practices",
            "150–300+ households, $100M–$500M+ AUM — advisors who live inside the plan",
            "Practices that need a tool they can trust in a client meeting, not a sales aid",
            "Not a consumer app. Not a robo. A professional instrument.",
        ],
    },
    {
        "title": "Where this fits for you",
        "bullets": [
            "Tax decisions don't exist in isolation — every Roth conversion, every bracket move affects the retirement trajectory",
            "Your platform likely shows the tax cost of a decision. Parallax shows the retirement cost.",
            "Integration: account balances, projected tax rates, Roth windows flow directly into the engine",
            "White-label: retirement planning layer inside your existing advisor workflow",
            "Distribution: your advisor or client base gains a best-in-class planning engine",
            "The combination nobody has built: tax + retirement in one coherent, honest model",
        ],
    },
    {
        "title": "Status + next steps",
        "bullets": [
            "Working prototype today — Scenarios and Sequencing tabs are live",
            "Engine: 98 years of data, full Monte Carlo, tax-aware account modeling",
            "This meeting: understand your platform, your data, your advisor workflow",
            "Next: live demo on a real household scenario",
            "Explore: integration, white-label, co-development, distribution, or investment",
            "",
            "Nathan Robinson, CFP   |   nathan.robinson007@gmail.com",
        ],
    },
]


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_slide(prs, slide_data):
    blank = prs.slide_layouts[6]   # completely blank
    slide = prs.slides.add_slide(blank)

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG

    title_text = slide_data["title"]
    bullets     = slide_data.get("bullets", [])
    is_cover    = slide_data.get("note") == "Cover"

    if is_cover:
        # Centred title + sub-bullets
        txb = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10.0), Inches(1.2))
        tf = txb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(52)
        run.font.bold = True
        run.font.color.rgb = INK
        run.font.name = "Calibri"

        txb2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.4), Inches(10.0), Inches(3.0))
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        for i, b in enumerate(bullets):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            p2.text = b
            p2.alignment = PP_ALIGN.LEFT
            p2.space_after = Pt(6)
            run2 = p2.runs[0]
            run2.font.size = Pt(18)
            run2.font.color.rgb = MUTED
            run2.font.name = "Calibri"
    else:
        # Title bar
        txb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(11.8), Inches(0.9))
        tf = txb.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        run = p.runs[0]
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = INK
        run.font.name = "Calibri"

        # Thin rule
        from pptx.util import Emu
        from pptx.oxml.ns import qn
        from lxml import etree
        line = slide.shapes.add_connector(
            1,   # MSO_CONNECTOR.STRAIGHT
            Inches(0.6), Inches(1.25),
            Inches(12.73), Inches(1.25)
        )
        line.line.color.rgb = RGBColor(0xDD, 0xD8, 0xD0)
        line.line.width = Pt(0.75)

        # Bullets
        txb2 = slide.shapes.add_textbox(Inches(0.6), Inches(1.45), Inches(12.0), Inches(5.7))
        tf2 = txb2.text_frame
        tf2.word_wrap = True
        for i, b in enumerate(bullets):
            p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            if b == "":
                p2.text = ""
                continue
            p2.text = "—  " + b
            p2.space_after = Pt(5)
            p2.space_before = Pt(2)
            run2 = p2.runs[0]
            run2.font.size = Pt(17)
            run2.font.color.rgb = INK
            run2.font.name = "Calibri"

    return slide


prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

for s in slides:
    add_slide(prs, s)

out = "/home/user/Parallax/Parallax-Pitch-Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
